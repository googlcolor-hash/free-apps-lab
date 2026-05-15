from __future__ import annotations

import importlib.util
import json
from collections import Counter
from pathlib import Path


ROOT = Path(__file__).resolve().parents[2]
DATA_DIR = ROOT / 'data'
OUT_DIR = ROOT / 'deliverables'
APPS_JSON = DATA_DIR / 'apps.json'
STATS_JSON = DATA_DIR / 'stats.json'
OUT_PPTX = OUT_DIR / 'free_apps_overview.pptx'


def _require_python_pptx():
    if importlib.util.find_spec('pptx') is None:
        raise SystemExit('WAIT: python-pptx is not installed. Install it before generating PPTX: py -3 -m pip install python-pptx')


def _load_apps() -> list[dict]:
    if not APPS_JSON.exists():
        raise SystemExit(f'WAIT: missing {APPS_JSON}')
    data = json.loads(APPS_JSON.read_text(encoding='utf-8'))
    if not isinstance(data, list) or not data:
        raise SystemExit('WAIT: apps.json is empty or invalid')
    return data


def _normalize_platforms(raw: object) -> list[str]:
    if isinstance(raw, list):
        return [str(x).strip() for x in raw if str(x).strip()]
    if raw is None:
        return []
    val = str(raw).strip()
    if not val:
        return []
    if ',' in val:
        return [p.strip() for p in val.split(',') if p.strip()]
    return [val]


def _stats_from_apps(apps: list[dict]) -> dict:
    categories = Counter()
    platforms = Counter()
    open_source = 0
    recommended = 0

    for app in apps:
        categories[str(app.get('section_l1') or 'Uncategorized')] += 1
        for p in _normalize_platforms(app.get('platforms')):
            platforms[p] += 1
        if bool(app.get('is_open_source')):
            open_source += 1
        if bool(app.get('is_recommended')):
            recommended += 1

    return {
        'total': len(apps),
        'open_source_true': open_source,
        'recommended_true': recommended,
        'top_categories': categories.most_common(5),
        'top_platforms': platforms.most_common(5),
    }


def _load_stats(apps: list[dict]) -> dict:
    if STATS_JSON.exists():
        try:
            s = json.loads(STATS_JSON.read_text(encoding='utf-8'))
            flags = s.get('flags', {})
            return {
                'total': int(flags.get('total', len(apps))),
                'open_source_true': int(flags.get('open_source_true', 0)),
                'recommended_true': int(flags.get('recommended_true', 0)),
                'top_categories': [
                    (i.get('category'), int(i.get('count', 0)))
                    for i in s.get('top_categories', [])[:5]
                ],
                'top_platforms': sorted(
                    [(k, int(v)) for k, v in (s.get('platforms', {}) or {}).items()],
                    key=lambda x: x[1],
                    reverse=True,
                )[:5],
            }
        except Exception:
            pass
    return _stats_from_apps(apps)


def _add_bullets(slide, lines: list[str]):
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for idx, line in enumerate(lines):
        p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        p.text = line


def main():
    _require_python_pptx()
    from pptx import Presentation

    apps = _load_apps()
    stats = _load_stats(apps)

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = 'Free Apps Intelligence Hub'
    slide.placeholders[1].text = 'Обзор каталога бесплатных приложений\nИсточник: awesome-free-apps'

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Структура базы'
    _add_bullets(slide, [
        'Единая схема данных для приложений (id, category, platforms, flags, source).',
        'Артефакты: JSON/CSV + интерактивный web-каталог + аналитические срезы.',
        'Навигация по категориям, платформам и open-source/recommended признакам.',
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Ключевые метрики'
    _add_bullets(slide, [
        f"Всего приложений: {stats['total']}",
        f"Open Source = true: {stats['open_source_true']}",
        f"Recommended = true: {stats['recommended_true']}",
        'Топ категорий: ' + '; '.join([f"{k} ({v})" for k, v in stats['top_categories']]) if stats['top_categories'] else 'Топ категорий: нет данных',
        'Топ платформ: ' + '; '.join([f"{k} ({v})" for k, v in stats['top_platforms']]) if stats['top_platforms'] else 'Топ платформ: нет данных',
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Use-cases'
    _add_bullets(slide, [
        'Быстрый подбор бесплатного инструмента под задачу команды.',
        'Сравнение категорий и платформ для стандартизации внутреннего стека.',
        'Формирование shortlist для пилотов без лицензионных затрат.',
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Roadmap'
    _add_bullets(slide, [
        'Автообновление данных по расписанию из upstream-репозитория.',
        'Скоринг приложений (активность, стабильность, community signals).',
        'Экспорт коротких рекомендаций под роли: dev/design/ops/marketing.',
    ])

    prs.save(OUT_PPTX)
    print(f'OK: created {OUT_PPTX}')


if __name__ == '__main__':
    main()
